"""Generate vivo cloud service analysis PPT (forest canopy style) - 24 slides
Updated: added competitive urgency, moat argument, KPIs, migration rules, unit economics"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Forest Canopy palette
BG       = RGBColor(0xFA, 0xF9, 0xF6)
BG2      = RGBColor(0xF0, 0xED, 0xE6)
TEXT     = RGBColor(0x2D, 0x3A, 0x2B)
TEXT2    = RGBColor(0x5A, 0x66, 0x56)
TEXT3    = RGBColor(0x8A, 0x95, 0x86)
OLIVE    = RGBColor(0x6B, 0x7A, 0x5E)
SAGE     = RGBColor(0xA8, 0xB0, 0x9C)
ACCENT   = RGBColor(0x3D, 0x6B, 0x4E)
WARM     = RGBColor(0x8B, 0x73, 0x55)
RED_S    = RGBColor(0xB0, 0x40, 0x30)
GREEN_BG = RGBColor(0xE8, 0xEE, 0xE4)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG  = RGBColor(0xFF, 0xFF, 0xFF)
BORDER   = RGBColor(0xC8, 0xC4, 0xB8)

SERIF = '宋体'
SANS  = '微软雅黑'
TOTAL = 26
IMG_DIR = r'D:\desktop\云服务截图'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def tb(slide, left, top, width, height, text, size=14, color=TEXT, bold=False,
       align=PP_ALIGN.LEFT, font=SANS, anchor=MSO_ANCHOR.TOP, spacing=1.5):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    try: tf.vertical_anchor = anchor
    except: pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    p.line_spacing = Pt(size * spacing)
    return txBox

def multi(slide, left, top, width, height, lines, anchor=MSO_ANCHOR.TOP):
    """lines: [(text, size, color, bold, align, font)]"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    try: tf.vertical_anchor = anchor
    except: pass
    for i, (text, sz, clr, bld, al, fn) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = fn
        p.alignment = al
        p.space_after = Pt(4)
        p.line_spacing = Pt(sz * 1.5)
    return txBox

def line(slide, left, top, width, color=OLIVE, h=0.03):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def card(slide, left, top, width, height, border_color=BORDER):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid(); s.fill.fore_color.rgb = CARD_BG
    s.line.color.rgb = border_color; s.line.width = Pt(1)
    s.adjustments[0] = 0.06
    return s

def accent_tag(slide, x, y, text, clr=ACCENT, w=1.5):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.38))
    s.fill.solid(); s.fill.fore_color.rgb = clr; s.line.fill.background()
    s.adjustments[0] = 0.4
    tb(slide, x, y+0.02, w, 0.34, text, 11, WHITE, True, PP_ALIGN.CENTER)

def label(slide, text, left=0.6, top=0.4):
    # colored pill badge
    w = max(1.2, len(text) * 0.18 + 0.4)
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(0.4))
    s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()
    s.adjustments[0] = 0.4
    tb(slide, left + 0.08, top + 0.02, w - 0.16, 0.36, text, 12, WHITE, True, PP_ALIGN.CENTER, SANS)

def page_num(slide, num):
    tb(slide, 11.8, 6.9, 1.2, 0.3, f'{num:02d} / {TOTAL}', 9, TEXT3, False, PP_ALIGN.RIGHT, SANS)

def dot(slide, x, y, clr):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.1), Inches(0.1))
    s.fill.solid(); s.fill.fore_color.rgb = clr; s.line.fill.background()

# ════════════════════════════════════════════════
# 01 Cover
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
label(slide, '竞品分析')
tb(slide, 1.5, 1.8, 10, 1.2, 'vivo 云服务', 44, TEXT, True, PP_ALIGN.LEFT, SERIF)
tb(slide, 1.5, 3.0, 10, 1.2, '增值体系优化分析', 44, TEXT, True, PP_ALIGN.LEFT, SERIF)
tb(slide, 1.5, 4.4, 8, 0.5, '紧迫性论证 · 护城河分析 · 量化KPI · 迁移规则 · 单位经济', 16, TEXT2, False, PP_ALIGN.LEFT)
line(slide, 1.5, 5.2, 0.8, ACCENT, 0.04)
for i, t in enumerate(['竞品分析', '护城河论证', '量化KPI', '单位经济']):
    x = 1.5 + i * 1.8
    accent_tag(slide, x, 5.6, t, ACCENT)
tb(slide, 1.5, 6.3, 5, 0.3, '森林冠层 · 2026', 11, TEXT3, False, PP_ALIGN.LEFT)
page_num(slide, 1)

# ════════════════════════════════════════════════
# 02 TOC
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
label(slide, '目录')
line(slide, 1.4, 1.2, 0.08, ACCENT)
tb(slide, 1.6, 1.05, 5, 0.5, '分析框架', 28, TEXT, True, PP_ALIGN.LEFT, SERIF)

toc = [
    ('第一幕：理解云服务', ACCENT, ['03 云服务的界定：存·传·用']),
    ('第二幕：参照对象', OLIVE, ['04 夸克演化路径', '05 夸克会员结构', '06 夸克的矛盾']),
    ('第三幕：自身审视', WARM, ['07 借鉴与避免', '08 vivo 现状盘点', '09 付费逻辑对比', '10 核心差距']),
    ('第四幕：紧迫性与护城河', RED_S, ['11 竞争紧迫性', '12-13 竞品实证（截图）', '14 护城河论证']),
    ('第五幕：破局路径', SAGE, ['15 三阶段总览', '16-17 第一阶段（含KPI）', '18-19 第二阶段（含KPI）', '20-21 第三阶段（含KPI）']),
    ('第六幕：落地保障', TEXT3, ['22 KPI对照', '23 迁移规则', '24 总结', '25 数据来源', '26 收束']),
]
positions = [(1.2, 1.7), (1.2, 3.4), (1.2, 5.1), (7.0, 1.7), (7.0, 3.4), (7.0, 5.1)]
for (title, clr, items), (x, y) in zip(toc, positions):
    card(slide, x, y, 5.5, 1.5, clr)
    tb(slide, x+0.3, y+0.15, 4.9, 0.3, title, 14, TEXT, True, PP_ALIGN.LEFT, SANS)
    text = '  '.join(items)
    tb(slide, x+0.3, y+0.55, 4.9, 0.8, text, 11, TEXT2, False, PP_ALIGN.LEFT, SANS)
page_num(slide, 2)

# ════════════════════════════════════════════════
# 03 云服务界定
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
label(slide, '背景')
line(slide, 1.4, 1.2, 0.08, ACCENT)
tb(slide, 1.6, 1.05, 5, 0.5, '云服务的界定', 28, TEXT, True, PP_ALIGN.LEFT, SERIF)

cols = [
    ('存', SAGE, '数据的持久化保存', '相册、文件、联系人、录音、笔记'),
    ('传', OLIVE, '数据的流动与同步', '跨设备访问、多端协同、分享与传输'),
    ('用', WARM, '数据的理解与增值', '搜索、整理、摘要、生成'),
]
for i, (char, clr, sub, desc) in enumerate(cols):
    x = 1.2 + i * 3.8
    card(slide, x, 1.9, 3.5, 2.8, clr)
    tb(slide, x, 2.3, 3.5, 0.6, char, 36, clr, False, PP_ALIGN.CENTER, SERIF)
    tb(slide, x, 2.9, 3.5, 0.35, sub, 13, TEXT, True, PP_ALIGN.CENTER)
    tb(slide, x, 3.35, 3.5, 0.7, desc, 12, TEXT2, False, PP_ALIGN.CENTER)

card(slide, 1.2, 5.1, 11.0, 1.2, OLIVE)
line(slide, 1.2, 5.1, 11.0, OLIVE, 0.04)
tb(slide, 1.6, 5.35, 10.2, 0.8, '当前大多数云服务停留在"存"和"传"，"用"这一层几乎是空白。这也是整篇分析的核心切入点。', 14, TEXT, False, PP_ALIGN.CENTER)
page_num(slide, 3)

# ════════════════════════════════════════════════
# 04 夸克演化
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
label(slide, '竞品画像')
line(slide, 1.4, 1.2, 0.08, ACCENT)
tb(slide, 1.6, 1.05, 8, 0.5, '夸克的演化路径', 28, TEXT, True, PP_ALIGN.LEFT, SERIF)

card(slide, 1.2, 1.8, 10.9, 0.9)
tb(slide, 1.5, 1.95, 10.3, 0.6, '从主打简洁的浏览器，逐步演化为集浏览器、网盘、AI 工具、内容平台于一体的"云服务大杂烩"。', 13, TEXT)

dims = [
    ('存储与传输', SAGE, ['免费 10GB，SVIP 最高 6TB', '不限速上传下载']),
    ('影音播放', OLIVE, ['倍速、投屏、后台播放', '社区内容共享']),
    ('文件处理', WARM, ['PDF 全家桶、格式互转', '在线解压、OCR']),
    ('AI 特权', ACCENT, ['AI 问答、修图、写作', '录音转写、视频总结']),
]
for i, (t, clr, items) in enumerate(dims):
    x = 1.2 + i * 2.85
    card(slide, x, 3.0, 2.6, 2.4)
    dot(slide, x+0.25, 3.2, clr)
    tb(slide, x+0.45, 3.15, 1.8, 0.3, t, 13, TEXT, True)
    text = '\n'.join(f'• {it}' for it in items)
    tb(slide, x+0.25, 3.6, 2.2, 1.5, text, 11, TEXT2)
page_num(slide, 4)

# ════════════════════════════════════════════════
# 05 夸克会员
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
label(slide, '竞品画像')
line(slide, 1.4, 1.2, 0.08, ACCENT)
tb(slide, 1.6, 1.05, 8, 0.5, '夸克的多会员结构', 28, TEXT, True, PP_ALIGN.LEFT, SERIF)

members = [
    ('网盘 VIP：高频存储型', SAGE, '学生和职场人，核心诉求是稳定存取。付费动机是降低使用成本，而非获得更多功能。'),
    ('AI 会员：效率提升型', OLIVE, '学生、创作者、高频办公人群。在意"我能做更多事"，而不是"我有更多空间"。'),
    ('扫描会员：高频轻工具型', WARM, '移动端即时处理。使用频率不一定高，但每次使用都"很关键"。'),
    ('小说会员：内容消费型', TEXT3, '为内容本身付费，不是为功能付费。必须独立成体系。'),
]
positions = [(1.2, 1.8, 5.3, 2.0), (1.2, 4.0, 5.3, 2.0), (7.0, 1.8, 5.3, 1.8), (7.0, 3.8, 5.3, 1.8)]
for (t, clr, desc), (x, y, w, h) in zip(members, positions):
    card(slide, x, y, w, h, clr)
    dot(slide, x+0.25, y+0.2, clr)
    tb(slide, x+0.45, y+0.15, w-0.6, 0.3, t, 13, TEXT, True)
    tb(slide, x+0.25, y+0.6, w-0.5, h-0.7, desc, 11, TEXT2)
page_num(slide, 5)

# ════════════════════════════════════════════════
# 06 夸克矛盾
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
label(slide, '竞品画像')
line(slide, 1.4, 1.2, 0.08, ACCENT)
tb(slide, 1.6, 1.05, 8, 0.5, '夸克的矛盾', 28, TEXT, True, PP_ALIGN.LEFT, SERIF)

card(slide, 1.2, 1.8, 10.9, 0.8, OLIVE)
tb(slide, 1.5, 1.9, 10.3, 0.6, '分层合理，但体验割裂。四种业务对应四种不同的使用频率、付费动机和成本结构。', 13, TEXT)

problems = [
    ('感知割裂', '"怎么这个功能又要开另一个会员？""我开的会员能用这个功能吗？"'),
    ('跨场景成本', 'AI 能力无法自然迁移到网盘或扫描场景'),
    ('一站式优势削弱', '使用分散的 AI 应用需要分别开会员；集中到夸克了还是需要各种会员'),
]
for i, (t, desc) in enumerate(problems):
    x = 1.2 + i * 3.7
    line(slide, x, 2.9, 3.3, RED_S, 0.04)
    tb(slide, x+0.1, 3.0, 3.2, 0.3, t, 14, RED_S, True)
    tb(slide, x+0.1, 3.4, 3.2, 1.2, desc, 12, TEXT2)

card(slide, 1.2, 5.0, 10.9, 1.0, OLIVE)
line(slide, 1.2, 5.0, 10.9, OLIVE, 0.04)
tb(slide, 1.5, 5.2, 10.3, 0.6, '夸克的矛盾不是"会员太多"，而是业务合理拆分与产品体验统一性之间的张力。', 14, TEXT, False, PP_ALIGN.CENTER)
page_num(slide, 6)

# ════════════════════════════════════════════════
# 07 借鉴与避免
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
label(slide, '战略判断')
line(slide, 1.4, 1.2, 0.08, ACCENT)
tb(slide, 1.6, 1.05, 8, 0.5, '从夸克到 vivo', 28, TEXT, True, PP_ALIGN.LEFT, SERIF)

card(slide, 1.2, 1.8, 5.3, 3.2, OLIVE)
multi(slide, 1.5, 2.0, 4.7, 2.8, [
    ('值得借鉴', 16, OLIVE, True, PP_ALIGN.LEFT, SANS),
    ('付费逻辑的转变：从"买空间"转向"买能力"。', 13, TEXT, False, PP_ALIGN.LEFT, SANS),
    ('用户愿意持续付费，不是因为空间更多，而是因为"我能做更多事"。', 12, TEXT2, False, PP_ALIGN.LEFT, SANS),
    ('用户分层逻辑也值得参考：不同场景的用户，付费动机根本不同。', 12, TEXT2, False, PP_ALIGN.LEFT, SANS),
])

tb(slide, 6.65, 3.0, 0.8, 0.4, 'VS', 13, TEXT3, True, PP_ALIGN.CENTER)

card(slide, 7.0, 1.8, 5.3, 3.2, RGBColor(0xD0, 0xA0, 0xA0))
multi(slide, 7.3, 2.0, 4.7, 2.8, [
    ('必须避免', 16, RED_S, True, PP_ALIGN.LEFT, SANS),
    ('第一，vivo 是手机厂商，不是内容平台。不应该靠版权内容做变现。', 13, TEXT, False, PP_ALIGN.LEFT, SANS),
    ('第二，功能堆砌会破坏一体化优势。vivo 应以"精简、一体化"为核心。', 12, TEXT2, False, PP_ALIGN.LEFT, SANS),
])

card(slide, 1.2, 5.4, 11.0, 1.0, OLIVE)
line(slide, 1.2, 5.4, 11.0, OLIVE, 0.04)
tb(slide, 1.5, 5.55, 10.4, 0.7, 'vivo 不应该学夸克做什么，而应该学夸克为什么用户愿意付费，然后用自己独有的方式去实现。', 14, TEXT, False, PP_ALIGN.CENTER)
page_num(slide, 7)

# ════════════════════════════════════════════════
# 08 vivo现状
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
label(slide, '自身现状')
line(slide, 1.4, 1.2, 0.08, ACCENT)
tb(slide, 1.6, 1.05, 8, 0.5, 'vivo 云服务现状', 28, TEXT, True, PP_ALIGN.LEFT, SERIF)

card(slide, 1.2, 1.8, 5.3, 3.0, SAGE)
dot(slide, 1.5, 2.0, SAGE)
tb(slide, 1.7, 1.95, 4.5, 0.3, '云服务核心功能', 14, TEXT, True)
items = ['相册/视频自动备份（20+ 种数据类型）', '我的云盘（类网盘）', '设备云备份与整机恢复',
         '微信/QQ 自动备份（VIP 专属）', '云端视频播放、云端解压（VIP 专属）',
         '免费 5GB；VIP 四档 50GB/200GB/1TB/2TB']
tb(slide, 1.5, 2.4, 4.8, 2.2, '\n'.join(f'• {it}' for it in items), 11, TEXT2)

card(slide, 7.0, 1.8, 5.3, 3.0, OLIVE)
dot(slide, 7.3, 2.0, OLIVE)
tb(slide, 7.5, 1.95, 4.5, 0.3, '跨设备协同（vivo 办公套件）', 14, TEXT, True)
items2 = ['原子笔记多端同步（含小V写作）', '手机与电脑剪贴板共享', '镜像投屏 / 远程控制',
          '文件互传（手机/平板/电脑）', '键鼠协同（接入 iPad/电脑）']
tb(slide, 7.3, 2.4, 4.8, 2.2, '\n'.join(f'• {it}' for it in items2), 11, TEXT2)

card(slide, 1.2, 5.1, 11.0, 1.5, OLIVE)
line(slide, 1.2, 5.1, 11.0, OLIVE, 0.04)
multi(slide, 1.5, 5.2, 10.4, 1.3, [
    ('AI 能力现状（来源：vivo 官网 / 2024 开发者大会）', 12, ACCENT, True, PP_ALIGN.LEFT, SANS),
    ('OriginOS 5 相册语义搜索（自然语言搜照片）· OriginOS 6 AI修图（消除/超清/扩图/调色）', 10, TEXT2, False, PP_ALIGN.LEFT, SANS),
    ('蓝心大模型 1B/3B/7B/13B/30B/千亿全矩阵 · 覆盖 5亿+ 用户 · 3万亿+ token 输出 · 端云协同架构已就绪（VCAP 平台）', 10, TEXT2, False, PP_ALIGN.LEFT, SANS),
    ('以上全部属于 OriginOS 系统级功能。vivo 云服务产品本身没有任何 AI 增值服务。', 11, RED_S, True, PP_ALIGN.LEFT, SANS),
])
page_num(slide, 8)

# ════════════════════════════════════════════════
# 09 付费对比
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
label(slide, '核心洞察')
line(slide, 1.4, 1.2, 0.08, ACCENT)
tb(slide, 1.6, 1.05, 8, 0.5, '付费逻辑的本质差异', 28, TEXT, True, PP_ALIGN.LEFT, SERIF)

card(slide, 1.2, 1.8, 5.0, 2.2, RGBColor(0xD0, 0xA0, 0xA0))
multi(slide, 1.5, 1.95, 4.4, 1.8, [
    ('vivo：买空间 + 基础特权', 15, RED_S, True, PP_ALIGN.LEFT, SANS),
    ('VIP 权益全部围绕存储与基础操作。用户付费后在功能层面没有任何获得感提升，续费动力弱。', 12, TEXT2, False, PP_ALIGN.LEFT, SANS),
])

tb(slide, 6.35, 2.5, 0.6, 0.4, 'VS', 12, TEXT3, True, PP_ALIGN.CENTER)

card(slide, 7.0, 1.8, 5.0, 2.2, OLIVE)
multi(slide, 7.3, 1.95, 4.4, 1.8, [
    ('夸克：买能力', 15, ACCENT, True, PP_ALIGN.LEFT, SANS),
    ('AI 能力、文件处理、影音特权均为会员专属。用户感知到"我能做更多事"，留存与续费逻辑更强。', 12, TEXT2, False, PP_ALIGN.LEFT, SANS),
])

rows = [
    ('', 'vivo 云服务', '夸克云盘'),
    ('免费空间', '5GB', '10GB'),
    ('最大容量', '2TB', '10TB'),
    ('付费逻辑', '买空间 + 基础特权', '买能力（AI/影音/文件处理）'),
    ('AI 增值', '无（AI 在端侧，不在云服务中）', '有（AI 问答/写作/转写）'),
    ('产品定位', '手机配套服务', '独立第三方网盘'),
]
for i, (c1, c2, c3) in enumerate(rows):
    y = 4.3 + i * 0.38
    clr = TEXT3 if i == 0 else TEXT2
    bld = i == 0
    sz = 10 if i == 0 else 11
    tb(slide, 1.2, y, 2.5, 0.35, c1, sz, clr, bld)
    tb(slide, 3.7, y, 4.0, 0.35, c2, sz, clr if i > 0 else TEXT3, bld)
    tb(slide, 7.7, y, 4.5, 0.35, c3, sz, ACCENT if i > 0 else TEXT3, i > 0)
    if i == 0:
        line(slide, 1.2, y + 0.33, 10.9, OLIVE, 0.02)
page_num(slide, 9)

# ════════════════════════════════════════════════
# 10 核心差距
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
label(slide, '核心洞察')
tb(slide, 3.5, 1.3, 6.3, 0.35, 'THE GAP', 11, ACCENT, True, PP_ALIGN.CENTER)
tb(slide, 2.5, 2.0, 8.3, 1.2, '端侧有 AI，云端没 AI', 32, TEXT, True, PP_ALIGN.CENTER, SERIF)

line(slide, 1.5, 3.5, 4.8, RED_S, 0.04)
multi(slide, 1.5, 3.65, 4.8, 1.8, [
    ('能力断层', 14, RED_S, True, PP_ALIGN.LEFT, SANS),
    ('蓝心大模型全矩阵已就绪（1B~千亿），端云协同架构（VCAP）已上线。', 11, TEXT2, False, PP_ALIGN.LEFT, SANS),
    ('OriginOS 5 相册语义搜索、AI 修图、蓝心小V 均为系统级功能。', 10, TEXT3, False, PP_ALIGN.LEFT, SANS),
    ('云服务产品本身无任何 AI 增值服务。', 11, RED_S, True, PP_ALIGN.LEFT, SANS),
])

line(slide, 7.0, 3.5, 4.8, RED_S, 0.04)
multi(slide, 7.0, 3.65, 4.8, 1.8, [
    ('数据孤立', 14, RED_S, True, PP_ALIGN.LEFT, SANS),
    ('照片、录音、笔记、联系人通过云服务同步到云端，但云端层面无横向连接。', 11, TEXT2, False, PP_ALIGN.LEFT, SANS),
    ('数据在云端，AI 却只在端侧。', 11, RED_S, True, PP_ALIGN.LEFT, SANS),
    ('官方策略："端侧优先，端云协同"——但云侧尚未执行。', 10, TEXT3, False, PP_ALIGN.LEFT, SANS),
])

card(slide, 2.0, 5.8, 9.3, 0.8, OLIVE)
line(slide, 2.0, 5.8, 9.3, OLIVE, 0.04)
tb(slide, 2.3, 5.9, 8.7, 0.5, '数据来源：vivo 2024 开发者大会 / vivo 官网帮助中心 / vivo 开放平台', 10, TEXT3, False, PP_ALIGN.CENTER)
page_num(slide, 10)

# ════════════════════════════════════════════════
# 11 竞争紧迫性（NEW）
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
label(slide, '紧迫性')
tb(slide, 3.5, 1.0, 6.3, 0.35, 'URGENCY', 11, RED_S, True, PP_ALIGN.CENTER)
tb(slide, 1.8, 1.5, 9.7, 0.8, '不做会怎样？', 30, TEXT, True, PP_ALIGN.CENTER, SERIF)
tb(slide, 2.0, 2.3, 9.3, 0.5, '云服务正在从"存储竞争"转向"AI 能力竞争"，窗口期不超过 12 个月。', 14, TEXT2, False, PP_ALIGN.CENTER)

# competitive table
comp = [
    ('厂商', 'AI 云服务进展', '用户规模', '状态'),
    ('Apple', 'Apple Intelligence + iCloud', '—', '海外已上线'),
    ('Google', 'Google One + Gemini', '15亿用户', '已上线'),
    ('夸克', 'AI超级框 + 网盘AI工具', '2亿用户 / 月活1.5亿', '已上线'),
    ('百度网盘', 'AI文档总结、图片语义搜索', '7亿+累计 / 月活1.5亿', '已上线'),
    ('小米', '云相册已内置AI搜索（人脸/场景/OCR）', '—', '已上线'),
    ('华为', '盘古大模型 + 华为云空间', '—', '部分上线'),
    ('荣耀', '魔法大模型已就绪，云服务无AI', '—', '能力闲置'),
    ('vivo', '蓝心大模型全矩阵就绪(1B~千亿)+VCAP端云协同，云服务无AI', '1.93亿台活跃设备', '能力闲置'),
]
for i, (c1, c2, c3, c4) in enumerate(comp):
    y = 2.9 + i * 0.30
    # header row = gray; vivo(8) and honor(7) = red
    is_header = i == 0
    is_vivo = i == 8
    is_honor = i == 7
    is_red = is_vivo or is_honor
    clr = TEXT3 if is_header else (RED_S if is_red else TEXT2)
    bld = is_header or is_vivo or is_honor
    tb(slide, 1.5, y, 1.6, 0.3, c1, 9, clr, bld)
    tb(slide, 3.1, y, 4.3, 0.3, c2, 9, clr, bld)
    tb(slide, 7.4, y, 3.0, 0.3, c3, 9, TEXT3 if is_header else (RED_S if is_red else TEXT2), is_red)
    tb(slide, 10.4, y, 1.8, 0.3, c4, 9, RED_S if is_red else clr, bld)
    if is_header:
        line(slide, 1.5, y + 0.28, 10.7, RED_S, 0.02)

card(slide, 1.2, 5.8, 11.0, 0.9, RED_S)
line(slide, 1.2, 5.8, 11.0, RED_S, 0.04)
multi(slide, 1.5, 5.85, 10.4, 0.75, [
    ('小米已在云相册落地AI搜索，而vivo与荣耀的AI能力仍停留在端侧，未与云服务打通。', 13, TEXT, False, PP_ALIGN.CENTER, SANS),
    ('数据来源：i.mi.com官方指南 | Google官方2025 | 新华网2025.01 | 百度官方/智研咨询 | QuestMobile 2025.06', 8, TEXT3, False, PP_ALIGN.CENTER, SANS),
])
page_num(slide, 11)

# ════════════════════════════════════════════════
# 12 国内竞品 AI 云服务实况（NEW - screenshots）
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
label(slide, '竞品实证')
tb(slide, 3.5, 1.0, 6.3, 0.35, 'EVIDENCE', 11, RED_S, True, PP_ALIGN.CENTER)
tb(slide, 1.5, 1.4, 10.3, 0.7, '国内竞品 AI 云服务已全面落地', 26, TEXT, True, PP_ALIGN.CENTER, SERIF)
tb(slide, 2.0, 2.1, 9.3, 0.4, '以下是各厂商官网/产品页的真实截图，证明 AI 云服务不再是概念，而是已上线的现实。', 12, TEXT2, False, PP_ALIGN.CENTER)

shots_cn = [
    ('百度网盘 · AI 工具箱', '文档转换、AI PPT、修图、OCR 等全套 AI 能力', 'baidu-pan-aitools.png', SAGE),
    ('华为 · 小艺 AI 修图', '"一句话搞定复杂图像需求"，消除路人、换背景', 'huawei-xiaoyi-ai-edit.png', OLIVE),
    ('小米 · 云服务 AI 搜索', '云相册已内置人脸识别、场景识别、文字识别', 'xiaomi-cloud-features.png', WARM),
    ('小米 · AI 智能搜索指南', '官方帮助文档明确记载 AI 搜索功能使用方法', 'xiaomi-cloud-guide.png', TEXT3),
]
for i, (title, desc, fname, clr) in enumerate(shots_cn):
    x = 0.8 + (i % 2) * 6.2
    y = 2.7 + (i // 2) * 2.5
    card(slide, x, y, 5.8, 2.3, clr)
    tb(slide, x + 0.15, y + 0.08, 5.5, 0.3, title, 12, TEXT, True)
    tb(slide, x + 0.15, y + 0.35, 5.5, 0.25, desc, 9, TEXT3)
    img_path = os.path.join(IMG_DIR, fname)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(x + 0.15), Inches(y + 0.6), Inches(5.5), Inches(1.6))
page_num(slide, 12)

# ════════════════════════════════════════════════
# 13 海外标杆 AI 能力（NEW - screenshots）
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
label(slide, '海外标杆')
tb(slide, 3.5, 1.0, 6.3, 0.35, 'BENCHMARK', 11, ACCENT, True, PP_ALIGN.CENTER)
tb(slide, 1.5, 1.4, 10.3, 0.7, 'Apple Intelligence 的终态形态', 26, TEXT, True, PP_ALIGN.CENTER, SERIF)
tb(slide, 2.0, 2.1, 9.3, 0.4, 'vivo 第三阶段"个人知识库"的对标对象。苹果是封闭的英文生态，vivo 有本土化优势。', 12, TEXT2, False, PP_ALIGN.CENTER)

shots_en = [
    ('Apple Intelligence · Siri 生产力', 'Siri 跨设备协同：MacBook / iPad / iPhone 联动，\n工作、学习、运动全场景 AI 助手', 'apple-intelligence-productivity.png'),
    ('Apple · 通信 / 图像 / 拍摄 AI', '智能通信过滤、AI 图像创作、相机智能识别——\nAI 深度嵌入系统级体验的标杆', 'apple-iphone-features.png'),
]
for i, (title, desc, fname) in enumerate(shots_en):
    x = 0.8 + i * 6.2
    card(slide, x, 2.7, 5.8, 4.2)
    tb(slide, x + 0.15, 2.8, 5.5, 0.3, title, 13, TEXT, True)
    tb(slide, x + 0.15, 3.15, 5.5, 0.45, desc, 10, TEXT3)
    img_path = os.path.join(IMG_DIR, fname)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(x + 0.15), Inches(3.7), Inches(5.5), Inches(3.1))

card(slide, 1.2, 6.5, 11.0, 0.5, OLIVE)
line(slide, 1.2, 6.5, 11.0, OLIVE, 0.04)
tb(slide, 1.5, 6.58, 10.4, 0.35, '截图来源：apple.com/apple-intelligence 官方产品页  |  截取时间：2026.06', 9, TEXT3, False, PP_ALIGN.CENTER)
page_num(slide, 13)

# ════════════════════════════════════════════════
# 14 护城河
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
label(slide, '护城河')
tb(slide, 3.5, 1.0, 6.3, 0.35, 'MOAT', 11, ACCENT, True, PP_ALIGN.CENTER)
tb(slide, 1.5, 1.5, 10.3, 0.8, '为什么用户不为免费 DeepSeek 而为 vivo 付费？', 26, TEXT, True, PP_ALIGN.CENTER, SERIF)

moats = [
    ('零搬运', SAGE, 'vivo 云 AI 直接在备份数据上工作。不需要上传、不需要手动选择文件、不需要打开第三方 App。第三方 AI 看不到你的云端备份。'),
    ('全量关联', OLIVE, '同时看到相册+笔记+录音+文档+联系人之间的关联。"那次日本旅行"同时牵出机票PDF、行程笔记、200张照片。任何第三方都无法实现。'),
    ('隐私闭环', WARM, '数据不发送给第三方。不是"把你的文件发给外部AI处理"，而是"在你的数据旁边跑AI"。2026年合规环境下是真实卖点。'),
]
for i, (title, clr, desc) in enumerate(moats):
    x = 1.2 + i * 3.8
    card(slide, x, 2.6, 3.5, 3.2, clr)
    line(slide, x+0.3, 2.9, 1.0, clr, 0.03)
    tb(slide, x+0.3, 3.0, 2.8, 0.35, title, 16, TEXT, True)
    tb(slide, x+0.3, 3.5, 2.8, 2.0, desc, 11, TEXT2)

# bottom highlight
card(slide, 1.2, 6.1, 11.0, 0.7, OLIVE)
line(slide, 1.2, 6.1, 11.0, OLIVE, 0.04)
tb(slide, 1.5, 6.22, 10.4, 0.45, 'vivo 云端沉积着夸克永远触达不到的数据：通话录音、联系人关系、系统笔记、相册元数据。', 13, TEXT, False, PP_ALIGN.CENTER)
page_num(slide, 14)

# ════════════════════════════════════════════════
# 15 三阶段总览
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
label(slide, '优化方案')
line(slide, 1.4, 1.2, 0.08, ACCENT)
tb(slide, 1.6, 1.05, 8, 0.5, '三阶段路径', 28, TEXT, True, PP_ALIGN.LEFT, SERIF)
tb(slide, 1.6, 1.7, 8, 0.4, '先用系统级私人数据建立体验壁垒，再让变现结构跟着体验走。每阶段有明确KPI，达标后才推进。', 11, TEXT2)

stages = [
    ('01', '相册语义搜索', '从拍照型用户切入，验证AI搜索付费意愿。KPI：AI渗透率≥15%，转化率≥4%（锚点：夸克19%渗透率）。', ['宝妈/家庭', '旅行党']),
    ('02', '跨模块数据联通', '打通通话录音、笔记、相册、联系人。KPI：跨模块活跃>12%，场景包购买率≥8%（锚点：Freemium 2-5%）。', ['职场', '学生', '家庭']),
    ('03', '个人知识库', '从"找东西"升级为"理解你"。KPI：对话频率≥2次/周，按量ARPU>订阅×1.1（锚点：百度文库+60%）。', ['重度用户']),
]
for i, (num, title, desc, tags) in enumerate(stages):
    x = 1.2 + i * 3.8
    card(slide, x, 2.3, 3.5, 3.8)
    tb(slide, x+0.3, 2.5, 2.8, 0.5, num, 32, OLIVE, False, PP_ALIGN.LEFT)
    tb(slide, x+0.3, 3.1, 2.8, 0.35, title, 15, TEXT, True)
    tb(slide, x+0.3, 3.55, 2.8, 1.8, desc, 11, TEXT2)
    for j, t in enumerate(tags):
        tx = x + 0.3 + j * 1.1
        accent_tag(slide, tx, 5.3, t, GREEN_BG if False else ACCENT, 0.95)
        # override to light bg with border
    if i < 2:
        line(slide, x + 3.5, 4.2, 0.3, OLIVE, 0.02)
page_num(slide, 15)

# ════════════════════════════════════════════════
# 16 第一阶段：场景 + 用户分群
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
label(slide, '第一阶段 · 短期')
line(slide, 1.4, 1.2, 0.08, ACCENT)
tb(slide, 1.6, 1.05, 10, 0.5, '云端长期记忆 AI', 28, TEXT, True, PP_ALIGN.LEFT, SERIF)

card(slide, 1.2, 1.8, 5.3, 1.3)
multi(slide, 1.5, 1.9, 4.7, 1.1, [
    ('手机是工作记忆，云端是长期记忆', 14, ACCENT, True, PP_ALIGN.LEFT, SANS),
    ('8万张照片在云端，手机只留1万张。端侧AI搜1万张，云端AI搜全部8万张。', 12, TEXT2, False, PP_ALIGN.LEFT, SANS),
])

# user segmentation
card(slide, 1.2, 3.3, 5.3, 2.8)
line(slide, 1.5, 3.5, 1.0, ACCENT, 0.03)
tb(slide, 1.5, 3.55, 3, 0.3, '目标用户优先级', 14, TEXT, True)

segs = [
    ('宝妈/家庭', '★★★', '"找宝宝第一次走路的照片"', '极高', SAGE),
    ('旅行党', '★★★', '"去年夏天海边那次"', '高', OLIVE),
    ('普通用户', '★★', '"有猫的""那张日落"', '中', WARM),
]
for i, (seg, pri, scenario, willing, clr) in enumerate(segs):
    y = 4.05 + i * 0.65
    card(slide, 1.5, y, 4.7, 0.55, clr)
    tb(slide, 1.7, y+0.05, 1.2, 0.4, seg, 11, TEXT, True)
    tb(slide, 2.9, y+0.05, 2.0, 0.4, scenario, 10, TEXT2)
    tb(slide, 5.0, y+0.05, 0.8, 0.4, willing, 10, clr, True)

# quant data
card(slide, 7.0, 1.8, 5.3, 4.3)
line(slide, 7.3, 2.0, 1.0, ACCENT, 0.03)
tb(slide, 7.3, 2.05, 4, 0.3, '量化场景数据', 14, TEXT, True)

data_rows = [
    ('照片搜索', '全球用户平均~2,800张(PhotoAiD)\n重度用户(宝妈/旅行)1-3万张', '54%用户觉得搜索困难(ResearchGate)', '云端全量语义索引'),
    ('旅行回忆', '年均2-3次旅行\n单次500-2000张(推算)', '按时间翻找，效率低', 'AI按"事件"自动聚合'),
    ('人物/事件', '宝妈拍摄量远超平均\n(推算)', '无法按"关系"聚合', 'AI识别人物+事件聚类'),
]
for i, (scene, data, current, ai_sol) in enumerate(data_rows):
    y = 2.5 + i * 1.15
    tb(slide, 7.3, y, 2.0, 0.3, scene, 12, TEXT, True)
    tb(slide, 7.3, y+0.25, 4.8, 0.25, data, 10, TEXT3)
    tb(slide, 7.3, y+0.65, 2.3, 0.25, f'现状：{current}', 10, RED_S)
    tb(slide, 9.6, y+0.65, 2.5, 0.25, f'AI：{ai_sol}', 10, ACCENT)

card(slide, 1.2, 6.3, 11.0, 0.5, OLIVE)
line(slide, 1.2, 6.3, 11.0, OLIVE, 0.04)
tb(slide, 1.5, 6.38, 10.4, 0.35, 'PhotoAiD/PetaPixel 2025 | ResearchGate 照片搜索行为研究  |  未标注数据为合理推算', 9, TEXT3, False, PP_ALIGN.CENTER)
page_num(slide, 16)

# ════════════════════════════════════════════════
# 17 第一阶段：KPI + 会员
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
label(slide, '第一阶段 · 短期')
line(slide, 1.4, 1.2, 0.08, ACCENT)
tb(slide, 1.6, 1.05, 8, 0.5, 'KPI + 会员结构', 28, TEXT, True, PP_ALIGN.LEFT, SERIF)

# KPI table
card(slide, 1.2, 1.8, 5.3, 2.6)
line(slide, 1.5, 2.0, 1.0, ACCENT, 0.03)
tb(slide, 1.5, 2.05, 3, 0.3, '成功标志（KPI）', 14, TEXT, True)

kpis = [
    ('AI 语义搜索渗透率', '≥ 15%', '夸克全量AI渗透率19%¹ 首次上线保守取80%'),
    ('月均搜索次数', '≥ 3 次/月', '夸克19%渗透率说明多数人不搜,是低频刚需³'),
    ('AI 能力包转化率', '≥ 4%', 'Freemium行业2-5%⁴ 百度网盘AI收入+120%⁵'),
    ('AI功能30日留存率', '≥ 40%', 'SaaS产品均值39%⁶ AI功能应略高'),
]
for i, (metric, target, note) in enumerate(kpis):
    y = 2.5 + i * 0.45
    tb(slide, 1.5, y, 2.8, 0.3, metric, 11, TEXT2)
    tb(slide, 4.3, y, 1.5, 0.3, target, 11, ACCENT, True)
    tb(slide, 1.5, y+0.2, 4.8, 0.2, note, 9, TEXT3)

# advance condition
card(slide, 1.2, 4.6, 5.3, 0.7, OLIVE)
line(slide, 1.2, 4.6, 5.3, OLIVE, 0.03)
tb(slide, 1.5, 4.7, 4.7, 0.45, '推进条件：KPI连续2月达标 + NPS>40 → 阶段二\n¹woshipm  ²SamMobile  ³新华网  ⁴Recurly  ⁵亿欧  ⁶Pendo', 10, TEXT, False, PP_ALIGN.LEFT)

# membership
card(slide, 7.0, 1.8, 5.3, 3.5)
line(slide, 7.3, 2.0, 1.0, OLIVE, 0.03)
tb(slide, 7.3, 2.05, 3, 0.3, '统一会员 + 能力分层', 14, TEXT, True)

tiers = [
    ('免费', '云服务底座', SAGE, ['相册备份', '文件同步', '联系人/日历同步']),
    ('付费', 'AI 能力包', OLIVE, ['相册语义搜索（云端）', '智能整理', '自动摘要']),
    ('按需', '高级存储包', WARM, ['大容量扩展', '极速恢复', '多副本备份']),
]
for i, (tier, name, clr, items) in enumerate(tiers):
    y = 2.5 + i * 0.9
    card(slide, 7.3, y, 4.7, 0.8, clr)
    tb(slide, 7.5, y+0.05, 0.8, 0.3, tier, 10, clr, True)
    tb(slide, 8.3, y+0.05, 1.5, 0.3, name, 11, TEXT, True)
    text = '  '.join(items)
    tb(slide, 7.5, y+0.35, 4.3, 0.3, text, 9, TEXT2)
page_num(slide, 17)

# ════════════════════════════════════════════════
# 18 第二阶段：场景
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
label(slide, '第二阶段 · 中期')
line(slide, 1.4, 1.2, 0.08, ACCENT)
tb(slide, 1.6, 1.05, 10, 0.5, '打通私人数据，建立护城河', 28, TEXT, True, PP_ALIGN.LEFT, SERIF)

card(slide, 1.2, 1.8, 10.9, 1.0)
multi(slide, 1.5, 1.9, 10.3, 0.8, [
    ('相册体验成熟后，单点能力的天花板很快会到。用户真实的生活场景从来不是孤立的——一次出差，涉及日历、通话录音、笔记、相册。', 12, TEXT2, False, PP_ALIGN.LEFT, SANS),
])

# quant data
card(slide, 1.2, 3.0, 10.9, 2.3)
line(slide, 1.5, 3.2, 1.0, ACCENT, 0.03)
tb(slide, 1.5, 3.25, 3, 0.3, '量化场景数据', 14, TEXT, True)

data2 = [
    ('出差整理', '月均1-2次，涉及4个模块', '手动切换3-4个App，15-20分钟', 'AI按"事件"自动聚合'),
    ('客户沟通', '销售年均200+通话', '通话录音与笔记分离', '搜"张三"→通话+笔记+合同'),
    ('学习资料', '每学期50+小时录音，100+篇笔记', '录音与笔记无法交叉检索', '跨模态语义搜索'),
]
for i, (scene, data, current, ai_sol) in enumerate(data2):
    y = 3.7 + i * 0.5
    tb(slide, 1.5, y, 1.5, 0.3, scene, 11, TEXT, True)
    tb(slide, 3.0, y, 2.5, 0.3, data, 10, TEXT3)
    tb(slide, 5.5, y, 3.0, 0.3, f'现状：{current}', 10, RED_S)
    tb(slide, 8.5, y, 3.5, 0.3, f'AI：{ai_sol}', 10, ACCENT)

# entries
card(slide, 1.2, 5.5, 5.3, 1.2, SAGE)
dot(slide, 1.5, 5.7, SAGE)
tb(slide, 1.7, 5.65, 4.5, 0.3, '蓝心小V 入口', 14, TEXT, True)
tb(slide, 1.5, 6.0, 4.8, 0.5, '→ "整理我和张三的沟通记录"\n→ "上周开会说了什么"', 11, TEXT2)

card(slide, 7.0, 5.5, 5.3, 1.2, OLIVE)
dot(slide, 7.3, 5.7, OLIVE)
tb(slide, 7.5, 5.65, 4.5, 0.3, 'vivo 浏览器入口', 14, TEXT, True)
tb(slide, 7.3, 6.0, 4.8, 0.5, '夸克搜公网，vivo 浏览器搜用户自己的云端数据', 11, TEXT2)
page_num(slide, 18)

# ════════════════════════════════════════════════
# 19 第二阶段：KPI + 权益
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
label(slide, '第二阶段 · 中期')
line(slide, 1.4, 1.2, 0.08, ACCENT)
tb(slide, 1.6, 1.05, 8, 0.5, 'KPI + 场景权益', 28, TEXT, True, PP_ALIGN.LEFT, SERIF)

# KPI
card(slide, 1.2, 1.8, 5.3, 2.4)
line(slide, 1.5, 2.0, 1.0, ACCENT, 0.03)
tb(slide, 1.5, 2.05, 3, 0.3, '成功标志（KPI）', 14, TEXT, True)

kpis2 = [
    ('跨模块搜索周活跃度', '≥ 12%（付费用户）', '夸克AI渗透19%×0.6³ 跨模块门槛更高'),
    ('非相册类数据备份增长', '连续3月环比 > 5%', '百度网盘AI DAU同比+90%是首年效应¹'),
    ('场景包独立购买率', '≥ 8%', 'Freemium 2-5%⁴ 场景包精准匹配痛点溢价'),
    ('用户跨模块使用率', '≥ 2个模块/活跃用户', '验证数据联通价值的最低有效值'),
]
for i, (metric, target, note) in enumerate(kpis2):
    y = 2.5 + i * 0.4
    tb(slide, 1.5, y, 2.8, 0.3, metric, 11, TEXT2)
    tb(slide, 4.3, y, 1.5, 0.3, target, 10, ACCENT, True)
    tb(slide, 1.5, y+0.18, 4.8, 0.2, note, 9, TEXT3)

card(slide, 1.2, 4.4, 5.3, 0.6, OLIVE)
line(slide, 1.2, 4.4, 5.3, OLIVE, 0.03)
tb(slide, 1.5, 4.48, 4.7, 0.4, '推进条件：跨模块活跃>20% + 非相册增长连续3月为正 → 阶段三\n¹亿欧  ²Pendo  ³woshipm  ⁴Recurly', 10, TEXT)

# scenes
scenes = [
    ('学', '学习场景', SAGE, ['录音转写 + AI会议总结', '笔记自动整理', '学习资料云归档']),
    ('工', '工作场景', OLIVE, ['合同扫描与识别', '多设备协同处理', '文档快速调取']),
    ('家', '家庭场景', WARM, ['相册智能整理', '旅行记忆生成', '家庭共享空间']),
]
for i, (char, title, clr, items) in enumerate(scenes):
    x = 7.0
    y = 1.8 + i * 1.5
    card(slide, x, y, 5.3, 1.3, clr)
    tb(slide, x+0.3, y+0.1, 0.5, 0.5, char, 22, clr, False, PP_ALIGN.CENTER, SERIF)
    tb(slide, x+0.8, y+0.1, 2, 0.3, title, 13, TEXT, True)
    text = '  '.join(items)
    tb(slide, x+0.8, y+0.45, 4, 0.7, text, 10, TEXT2)
page_num(slide, 19)

# ════════════════════════════════════════════════
# 20 第三阶段：知识库
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
label(slide, '第三阶段 · 长期')
tb(slide, 3.5, 1.2, 6.3, 0.35, 'ENDGAME', 11, ACCENT, True, PP_ALIGN.CENTER)
tb(slide, 2.0, 1.8, 9.3, 0.8, '从"找东西"到"理解你"', 28, TEXT, True, PP_ALIGN.CENTER, SERIF)
tb(slide, 2.5, 2.8, 8.3, 0.5, '当数据联通和 AI 理解能力成熟后，vivo 云服务从工具升级为"懂你的私人 AI"。', 13, TEXT2, False, PP_ALIGN.CENTER)

card(slide, 2.0, 3.5, 4.3, 2.0)
multi(slide, 2.3, 3.65, 3.7, 1.6, [
    ('对话示例', 14, TEXT, True, PP_ALIGN.LEFT, SANS),
    ('→ "我最近在关注什么"', 12, TEXT, False, PP_ALIGN.LEFT, SANS),
    ('→ "整理和这个客户的往来"', 12, TEXT, False, PP_ALIGN.LEFT, SANS),
    ('→ "我去年这时候在做什么"', 12, TEXT, False, PP_ALIGN.LEFT, SANS),
    ('→ "帮我回顾这一年"', 12, TEXT, False, PP_ALIGN.LEFT, SANS),
])

card(slide, 7.0, 3.5, 4.3, 2.0)
line(slide, 7.3, 3.7, 1.0, ACCENT, 0.03)
tb(slide, 7.3, 3.75, 3, 0.3, '成功标志（KPI）', 14, TEXT, True)
kpis3 = [
    ('对话频率', '≥ 2次/周', '夸克大学生AI渗透80%是高频群体¹ 一般用户2次形成习惯'),
    ('按量ARPU', '> 订阅ARPU × 1.1', '百度文库付费率+60%验证AI拉升ARPU² 保守取1.1'),
    ('6个月留存', '≥ 55%', 'SaaS月流失1-7%³ AI留存提升20%+⁴ 取保守值'),
    ('迁移阻力', '自评 ≥ 7/10', '验证生态锁定效果的定性指标'),
]
for i, (m, t, n) in enumerate(kpis3):
    y = 4.15 + i * 0.35
    tb(slide, 7.3, y, 1.8, 0.25, m, 11, TEXT2)
    tb(slide, 9.1, y, 2.0, 0.25, t, 11, ACCENT, True)
    tb(slide, 7.3, y+0.18, 4.0, 0.15, n, 8, TEXT3)

card(slide, 1.5, 5.8, 10.3, 0.8, OLIVE)
line(slide, 1.5, 5.8, 10.3, OLIVE, 0.04)
multi(slide, 1.8, 5.88, 9.7, 0.6, [
    ('对标 Apple Intelligence 终态形态。苹果是封闭的英文生态，vivo 有本土化优势。', 12, TEXT, False, PP_ALIGN.CENTER, SANS),
    ('¹新华网  ²证券时报  ³CustomerGauge  ⁴SaaStr', 8, TEXT3, False, PP_ALIGN.CENTER, SANS),
])
page_num(slide, 20)

# ════════════════════════════════════════════════
# 21 第三阶段：计费
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
label(slide, '第三阶段 · 长期')
line(slide, 1.4, 1.2, 0.08, ACCENT)
tb(slide, 1.6, 1.05, 8, 0.5, '基础免费 + AI 按量计费', 28, TEXT, True, PP_ALIGN.LEFT, SERIF)
tb(slide, 1.6, 1.7, 8, 0.35, '当个人知识库体验成熟后，将 AI 能力拆分为可消耗型服务。', 12, TEXT2)

pricing = [
    ('免费', '基础服务', TEXT3, ['云存储作为基础设施免费提供', '基础备份与同步能力不收费']),
    ('按需付费', 'AI 能力（credit）', OLIVE, ['图片语义搜索调用次数', '录音转写分钟数', '文档总结任务', '智能整理与标签生成']),
    ('订阅制', '高级存储', WARM, ['大容量扩展', '高速上传下载通道', '多设备同步加速']),
]
for i, (tier, name, clr, items) in enumerate(pricing):
    x = 1.2 + i * 3.8
    card(slide, x, 2.2, 3.5, 3.0, clr)
    tb(slide, x+0.3, 2.35, 2.8, 0.25, tier, 10, clr, True)
    tb(slide, x+0.3, 2.65, 2.8, 0.3, name, 15, TEXT, True)
    text = '\n'.join(f'• {it}' for it in items)
    tb(slide, x+0.3, 3.1, 2.8, 1.5, text, 11, TEXT2)

card(slide, 1.2, 5.5, 11.0, 0.8, OLIVE)
line(slide, 1.2, 5.5, 11.0, OLIVE, 0.04)
tb(slide, 1.5, 5.62, 10.4, 0.5, '更贴近真实成本（AI 为消耗型资源），对低频用户更友好，避免了会员叠加带来的复杂感。', 13, TEXT, False, PP_ALIGN.CENTER)
page_num(slide, 21)

# ════════════════════════════════════════════════
# 22 三阶段KPI总览
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
label(slide, '量化框架')
line(slide, 1.4, 1.2, 0.08, ACCENT)
tb(slide, 1.6, 1.05, 8, 0.5, '三阶段 KPI 对照', 28, TEXT, True, PP_ALIGN.LEFT, SERIF)

# header
header = [('阶段', 1.5, 1.5), ('核心指标', 3.0, 3.0), ('达标线', 6.0, 2.5), ('推进条件', 8.5, 4.0)]
for text, x, w in header:
    tb(slide, x, 1.8, w, 0.3, text, 10, TEXT3, True)
line(slide, 1.5, 2.1, 10.5, OLIVE, 0.02)

kpi_all = [
    ('阶段一', SAGE, 'AI语义搜索渗透率', '≥ 15%', '转化率 ≥ 4%  锚点:夸克19%'),
    ('', SAGE, 'AI功能30日留存率', '≥ 40%', 'KPI连续2月达标+NPS>40'),
    ('阶段二', OLIVE, '跨模块搜索周活跃度', '≥ 12%', '场景包购买率 ≥ 8%  锚点:Freemium 2-5%'),
    ('', OLIVE, '非相册备份增长', '环比 > 5%', '活跃>20%+增长连续3月为正'),
    ('阶段三', WARM, '用户对话频率', '≥ 2次/周', '按量ARPU > 订阅 × 1.1  锚点:百度+60%'),
    ('', WARM, '6个月留存率', '≥ 55%', '迁移阻力自评 ≥ 7/10'),
]
for i, (stage, clr, metric, target, advance) in enumerate(kpi_all):
    y = 2.3 + i * 0.45
    tb(slide, 1.5, y, 1.5, 0.3, stage, 11, clr, bool(stage), PP_ALIGN.LEFT)
    tb(slide, 3.0, y, 3.0, 0.3, metric, 11, TEXT2)
    tb(slide, 6.0, y, 2.5, 0.3, target, 11, ACCENT, True)
    tb(slide, 8.5, y, 4.0, 0.3, advance, 10, TEXT3)
page_num(slide, 22)

# ════════════════════════════════════════════════
# 23 迁移规则
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
label(slide, '落地保障')
line(slide, 1.4, 1.2, 0.08, ACCENT)
tb(slide, 1.6, 1.05, 8, 0.5, '阶段间迁移规则', 28, TEXT, True, PP_ALIGN.LEFT, SERIF)
tb(slide, 1.6, 1.7, 8, 0.3, '核心原则：老用户权益只升不降，过渡无感。', 12, ACCENT, True)

# stage 1 -> 2
card(slide, 1.2, 2.2, 5.3, 2.4, SAGE)
line(slide, 1.5, 2.4, 1.2, SAGE, 0.03)
tb(slide, 1.5, 2.45, 3, 0.3, '阶段一 → 阶段二', 13, TEXT, True)
tb(slide, 1.5, 2.8, 4.7, 0.25, '统一会员 → 场景能力包', 11, SAGE, True)
multi(slide, 1.5, 3.1, 4.7, 1.3, [
    ('统一会员用户：自动升级为"全场景通行证"，价格不变，包含所有场景包。', 11, TEXT2, False, PP_ALIGN.LEFT, SANS),
    ('新用户：可选择"全场景通行证"或按场景单独购买。', 11, TEXT2, False, PP_ALIGN.LEFT, SANS),
])

# stage 2 -> 3
card(slide, 7.0, 2.2, 5.3, 2.4, OLIVE)
line(slide, 7.3, 2.4, 1.2, OLIVE, 0.03)
tb(slide, 7.3, 2.45, 3, 0.3, '阶段二 → 阶段三', 13, TEXT, True)
tb(slide, 7.3, 2.8, 4.7, 0.25, '场景能力包 → Credit 机制', 11, OLIVE, True)
multi(slide, 7.3, 3.1, 4.7, 1.3, [
    ('按历史3个月平均使用量折算为等额月度credit，月费不变。', 11, TEXT2, False, PP_ALIGN.LEFT, SANS),
    ('换算透明：设置页提供"我的AI使用报告"，用户可自查。', 11, TEXT2, False, PP_ALIGN.LEFT, SANS),
])

# bottom principles
card(slide, 1.2, 4.9, 11.0, 1.6, OLIVE)
line(slide, 1.2, 4.9, 11.0, OLIVE, 0.04)
tb(slide, 1.5, 5.05, 3, 0.3, '通用原则', 14, TEXT, True)
principles = [
    ('过渡期 ≥ 3个月', '新阶段上线后，老方案继续运行3个月'),
    ('主动通知+教育引导', 'App内弹窗+短信告知权益变化'),
    ('退出机制', '不满意可选择"保持原方案"或"全额退款"'),
    ('灰度放量', '先对5%用户开放，收集反馈后逐步扩大'),
]
for i, (p, desc) in enumerate(principles):
    x = 1.5 + (i % 2) * 5.3
    y = 5.4 + (i // 2) * 0.5
    tb(slide, x, y, 2.2, 0.3, f'• {p}', 11, TEXT, True)
    tb(slide, x+2.2, y, 3.0, 0.3, desc, 10, TEXT3)
page_num(slide, 25)

# ════════════════════════════════════════════════
# 24 总结
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
label(slide, '总结')
line(slide, 1.4, 1.2, 0.08, ACCENT)
tb(slide, 1.6, 1.05, 8, 0.5, '三个必须回答的问题', 28, TEXT, True, PP_ALIGN.LEFT, SERIF)

questions = [
    ('01', '为什么必须做？', RED_S,
     '竞对已在行动，12个月内不做就会被甩开。\n夸克AI渗透率19%、百度网盘AI收入+120%、Samsung S25 AI使用率70%。\n这不是增长机会，是生存题。'),
    ('02', '为什么 vivo 能做别人做不了？', ACCENT,
     '零搬运：数据已在云端，无需手动上传。\n全量关联：跨模块数据理解，第三方永远做不到。\n隐私闭环：数据不出基础设施，合规优势。'),
    ('03', '做了能赚多少钱？', OLIVE,
     'AI推理成本3年降280倍，语义索引为一次性成本。\n百度文库AI推动付费率+60%，Freemium转化率2-5%即可盈利。\n1.93亿台设备，结构性优势明确。'),
]
for i, (num, title, clr, desc) in enumerate(questions):
    y = 1.8 + i * 1.6
    card(slide, 1.2, y, 11.0, 1.4, clr)
    line(slide, 1.2, y, 11.0, clr, 0.04)
    tb(slide, 1.5, y+0.1, 0.6, 0.4, num, 24, clr, False, PP_ALIGN.LEFT)
    tb(slide, 2.2, y+0.1, 4, 0.3, title, 16, TEXT, True)
    tb(slide, 2.2, y+0.5, 9.5, 0.8, desc, 11, TEXT2)
page_num(slide, 25)

# ════════════════════════════════════════════════
# 25 数据来源
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
label(slide, '附录')
line(slide, 1.4, 1.2, 0.08, ACCENT)
tb(slide, 1.6, 1.05, 8, 0.5, '数据来源', 28, TEXT, True, PP_ALIGN.LEFT, SERIF)

sources_left = [
    ('¹ 夸克全量用户AI渗透率 19%', 'woshipm.com 2025', 'https://www.woshipm.com/ai/6293627.html'),
    ('² Samsung S25 Galaxy AI使用率 70%', 'SamMobile 2025', 'https://www.sammobile.com/news/galaxy-ai-adoption-soars-70-percent-of-galaxy-s25-users-are-on-board/'),
    ('³ 夸克大学生AI渗透率 80%+', '新华网 2025.09', 'http://www.news.cn/tech/20250910/c7b0db4619364261a2c0c37dd8d4a406/c.html'),
    ('⁴ Freemium转付费典型 2-5%', 'Recurly 2025', 'https://recurly.com/blog/what-is-freemium-a-guide-for-subscription-businesses/'),
    ('⁵ 百度网盘 AI DAU同比+90% / AI收入+120%', '亿欧 2025.01', 'https://www.iyiou.com/news/202601241120584'),
]
for i, (desc, source, url) in enumerate(sources_left):
    y = 1.7 + i * 0.7
    card(slide, 1.2, y, 5.3, 0.6)
    tb(slide, 1.5, y+0.02, 5.0, 0.25, desc, 10, TEXT, True)
    tb(slide, 1.5, y+0.22, 5.0, 0.18, f'{source}  |  {url}', 7, TEXT3)

sources_right = [
    ('⁶ SaaS产品1个月留存 39%', 'Pendo 2025', 'https://www.pendo.io/pendo-blog/user-retention-rate-benchmarks/'),
    ('⁷ 百度文库 付费率同比+60%', '证券时报 2025', 'https://www.stcn.com/article/detail/1504654.html'),
    ('⁸ AI用户留存比非AI高20%+', 'SaaStr 2025', 'https://www.saastr.com/if-your-ai-feature-didnt-materially-boost-revenue-it-doesnt-count-try-again/'),
    ('⁹ 夸克 月活同比+53% / 日活+75%', 'Xsignal 2025', 'https://www.jwview.com/jingwei/html/m/04-09/620307.shtml'),
    ('¹⁰ AI推理成本3年下降280倍', 'Stanford HAI 2025', 'https://hai.stanford.edu/ai-index/2025-ai-index-report'),
]
for i, (desc, source, url) in enumerate(sources_right):
    y = 1.7 + i * 0.7
    card(slide, 7.0, y, 5.3, 0.6)
    tb(slide, 7.3, y+0.02, 5.0, 0.25, desc, 10, TEXT, True)
    tb(slide, 7.3, y+0.22, 5.0, 0.18, f'{source}  |  {url}', 7, TEXT3)

card(slide, 1.2, 5.4, 11.0, 0.5, OLIVE)
line(slide, 1.2, 5.4, 11.0, OLIVE, 0.04)
tb(slide, 1.5, 5.52, 10.4, 0.3, '艾媒咨询《2025年中国个人网盘市场发展状况》: https://www.iimedia.cn/c400/106309.html', 7, TEXT3, False, PP_ALIGN.CENTER)

tb(slide, 1.5, 6.1, 10.4, 0.3, '未标注来源的数据为基于公开信息的合理推算，不代表任何机构官方数据。', 9, TEXT3, False, PP_ALIGN.CENTER)
page_num(slide, 25)

# ════════════════════════════════════════════════
# 26 收束
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
label(slide, '收束')
tb(slide, 3.5, 1.2, 6.3, 0.35, 'CLOSING', 11, ACCENT, True, PP_ALIGN.CENTER)
page_num(slide, 24)
tb(slide, 2.0, 1.8, 9.3, 1.2, '把分散的私人数据变成一个\n可对话的个人知识库', 30, TEXT, True, PP_ALIGN.CENTER, SERIF)

nodes = ['蓝心小V', 'vivo 浏览器', '个人知识库']
for i, n in enumerate(nodes):
    x = 2.5 + i * 3.0
    accent_tag(slide, x, 3.4, n, ACCENT, 1.8)
    if i < 2:
        tb(slide, x + 1.9, 3.42, 0.5, 0.36, '+' if i == 0 else '→', 16, TEXT3, False, PP_ALIGN.CENTER)

multi(slide, 2.5, 4.1, 8.3, 1.5, [
    ('入口层：蓝心小V + vivo 浏览器', 14, TEXT, False, PP_ALIGN.CENTER, SANS),
    ('内容层：云端私人数据', 14, TEXT, False, PP_ALIGN.CENTER, SANS),
    ('理解层：蓝心大模型', 14, TEXT, False, PP_ALIGN.CENTER, SANS),
    ('三者打通，形成夸克无法复制的差异化能力。', 13, TEXT2, False, PP_ALIGN.CENTER, SANS),
])

line(slide, 6.2, 5.8, 0.9, ACCENT, 0.03)
tb(slide, 2.5, 5.9, 8.3, 0.4, 'vivo 的优势不在于功能比夸克更多，而在于数据比任何第三方都更深', 12, TEXT3, False, PP_ALIGN.CENTER)

# Save
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'output', 'vivo-cloud-analysis-v3.pptx')
prs.save(out)
print(f'Saved: {out}')
