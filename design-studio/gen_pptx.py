"""Generate vivo cloud service competitive analysis PPT (zen garden style)"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
BG       = RGBColor(0xF8, 0xFA, 0xF8)
BG2      = RGBColor(0xF0, 0xF4, 0xF0)
TEXT     = RGBColor(0x2A, 0x30, 0x2A)
TEXT2    = RGBColor(0x6A, 0x7A, 0x6A)
TEXT3    = RGBColor(0x9A, 0xAA, 0x9A)
CYAN     = RGBColor(0x7D, 0xCE, 0xDC)
GREEN    = RGBColor(0xAD, 0xDD, 0xC0)
INDIGO   = RGBColor(0xC1, 0xD0, 0xFF)
RED_S    = RGBColor(0xC8, 0x70, 0x70)
WARM     = RGBColor(0xA0, 0x80, 0x60)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

def set_slide_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=TEXT, bold=False, align=PP_ALIGN.LEFT, font_name='微软雅黑',
                 anchor=MSO_ANCHOR.TOP, line_spacing=1.5):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    p.space_after = Pt(0)
    p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_multi_text(slide, left, top, width, height, lines, anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, size, color, bold, align, font_name)"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    try:
        tf.vertical_anchor = anchor
    except:
        pass
    for i, (text, size, color, bold, align, fname) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = fname
        p.alignment = align
        p.space_after = Pt(4)
        p.line_spacing = Pt(size * 1.6)
    return txBox

def add_line(slide, left, top, width, color=CYAN, height=0.02):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, border_color=RGBColor(0xD0,0xDA,0xC8)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF,0xFF,0xFF)
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    shape.adjustments[0] = 0.05
    return shape

def add_label(slide, text, left=0.5, top=0.4):
    add_text_box(slide, left, top, 3, 0.4, text, 10, TEXT3, False, PP_ALIGN.LEFT, '微软雅黑')

def add_page_num(slide, num, total=10):
    add_text_box(slide, 11.5, 6.8, 1.5, 0.4, f'{num:02d} / {total}', 10, TEXT3, False, PP_ALIGN.RIGHT)

# ═══════════════════════════════════════
# 01 Cover
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_label(slide, '竞品分析')
add_text_box(slide, 1.5, 2.2, 10.3, 1.5, 'vivo 云服务 × 夸克云盘',
             42, TEXT, True, PP_ALIGN.CENTER, '微软雅黑', MSO_ANCHOR.MIDDLE)
add_text_box(slide, 2, 3.6, 9.3, 0.8, '识别 vivo 云服务的改进空间',
             18, TEXT2, False, PP_ALIGN.CENTER)
add_line(slide, 6.1, 4.5, 1)
add_text_box(slide, 3, 5.0, 7.3, 0.5, '枯山水 · 竞品研究 · 2025',
             13, TEXT3, False, PP_ALIGN.CENTER)
add_page_num(slide, 1)

# ═══════════════════════════════════════
# 02 云服务界定
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_label(slide, '背景')
add_line(slide, 1.2, 1.1, 0.08)
add_text_box(slide, 1.5, 0.9, 5, 0.7, '云服务的界定', 26, TEXT, True)

add_card(slide, 3.5, 2.2, 6.3, 2.0)
add_text_box(slide, 3.8, 2.5, 5.7, 0.6, '存 · 传 · 用', 28, TEXT, True, PP_ALIGN.CENTER, '微软雅黑')
add_text_box(slide, 3.8, 3.3, 5.7, 0.7, '手机厂商与三方网盘提供的云服务\n核心围绕用户的数字资产展开', 14, TEXT2, False, PP_ALIGN.CENTER)

add_card(slide, 2.5, 4.8, 8.3, 1.5, CYAN)
add_multi_text(slide, 2.8, 5.0, 7.7, 1.2, [
    ('分析对象：vivo 云服务 vs 夸克云盘', 14, TEXT, True, PP_ALIGN.CENTER, '微软雅黑'),
    ('一个代表手机厂商的系统级云服务，一个代表三方网盘的能力堆叠——两者走出了截然不同的路径。', 13, TEXT2, False, PP_ALIGN.CENTER, '微软雅黑'),
])
add_page_num(slide, 2)

# ═══════════════════════════════════════
# 03 夸克演化
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_label(slide, '竞品画像')
add_line(slide, 1.2, 1.1, 0.08)
add_text_box(slide, 1.5, 0.9, 5, 0.7, '夸克的演化路径', 26, TEXT, True)

add_card(slide, 1.2, 1.8, 10.9, 1.0)
add_text_box(slide, 1.5, 1.9, 10.3, 0.7, '从主打简洁的浏览器，逐渐发展为云服务"大杂烩"。功能越堆越多，饱受一部分用户诟病。',
             14, TEXT)

# 3 columns
cols = [
    ('存储与传输', CYAN, ['免费 10GB，SVIP 最高 6TB', '不限速上传下载（会员）']),
    ('影音播放', GREEN, ['倍速 / 后台 / 投屏', '社区内容共享']),
    ('文件处理', INDIGO, ['PDF 全家桶（互转、OCR）', '文件在线解压']),
]
for i, (title, accent, items) in enumerate(cols):
    x = 1.2 + i * 3.7
    add_card(slide, x, 3.1, 3.4, 2.0)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.3), Inches(3.3), Inches(0.12), Inches(0.12))
    shape.fill.solid(); shape.fill.fore_color.rgb = accent; shape.line.fill.background()
    add_text_box(slide, x + 0.5, 3.25, 2.5, 0.4, title, 14, TEXT, True)
    text = '\n'.join(f'• {it}' for it in items)
    add_text_box(slide, x + 0.3, 3.8, 3.0, 1.2, text, 12, TEXT2)

# AI card
add_card(slide, 1.2, 5.4, 10.9, 1.5)
shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), Inches(5.6), Inches(0.12), Inches(0.12))
shape.fill.solid(); shape.fill.fore_color.rgb = WARM; shape.line.fill.background()
add_text_box(slide, 1.7, 5.55, 2, 0.4, 'AI 特权', 14, TEXT, True)
ai_items = ['AI 问答（通义大模型）', 'AI 修图、写作、解题', '录音转写（1200 分钟/月）', 'AI 字幕（视频自动生成）']
add_text_box(slide, 1.5, 6.0, 4.8, 0.9, '\n'.join(f'• {it}' for it in ai_items[:2]), 12, TEXT2)
add_text_box(slide, 6.5, 6.0, 4.8, 0.9, '\n'.join(f'• {it}' for it in ai_items[2:]), 12, TEXT2)
add_page_num(slide, 3)

# ═══════════════════════════════════════
# 04 vivo 云服务
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_label(slide, '自身现状')
add_line(slide, 1.2, 1.1, 0.08)
add_text_box(slide, 1.5, 0.9, 8, 0.7, 'vivo 云服务：两条产品线', 26, TEXT, True)

# Two columns
add_card(slide, 1.2, 1.8, 5.5, 3.8)
shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), Inches(2.0), Inches(0.12), Inches(0.12))
shape.fill.solid(); shape.fill.fore_color.rgb = CYAN; shape.line.fill.background()
add_text_box(slide, 1.7, 1.95, 4, 0.4, '云端备份', 14, TEXT, True)
backup_items = ['相册自动备份（照片/视频）', '联系人、日历、录音机同步', '云盘文件上传与管理', '回收站（最近删除）', '免费 5GB；付费最高 2048GB']
add_text_box(slide, 1.5, 2.5, 5.0, 3.0, '\n'.join(f'• {it}' for it in backup_items), 12, TEXT2)

add_card(slide, 7.0, 1.8, 5.5, 3.8)
shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.3), Inches(2.0), Inches(0.12), Inches(0.12))
shape.fill.solid(); shape.fill.fore_color.rgb = INDIGO; shape.line.fill.background()
add_text_box(slide, 7.5, 1.95, 4, 0.4, '跨设备协同', 14, TEXT, True)
sync_items = ['原子笔记多端同步', '手机 ↔ 电脑剪贴板共享', '镜像投屏 / 远程控制', '文件互传（手机/平板/电脑）', '浏览器书签同步', '键鼠协同（接入 iPad/电脑）']
add_text_box(slide, 7.3, 2.5, 5.0, 3.0, '\n'.join(f'• {it}' for it in sync_items), 12, TEXT2)

# AI highlight
add_card(slide, 1.2, 5.9, 11.3, 1.1, CYAN)
add_multi_text(slide, 1.5, 6.0, 10.7, 0.9, [
    ('AI 现状：', 13, TEXT, True, PP_ALIGN.LEFT, '微软雅黑'),
    ('蓝心大模型已落地系统层（蓝心小V、通话摘要、图片语义搜索、笔记写作辅助），但尚未与云服务体系打通。', 13, TEXT2, False, PP_ALIGN.LEFT, '微软雅黑'),
], anchor=MSO_ANCHOR.MIDDLE)
add_page_num(slide, 4)

# ═══════════════════════════════════════
# 05 夸克会员体系
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_label(slide, '增值服务')
add_line(slide, 1.2, 1.1, 0.08)
add_text_box(slide, 1.5, 0.9, 8, 0.7, '夸克：分散的会员体系', 26, TEXT, True)

# Left column
add_card(slide, 1.2, 1.8, 5.5, 2.2)
shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), Inches(2.0), Inches(0.12), Inches(0.12))
shape.fill.solid(); shape.fill.fore_color.rgb = CYAN; shape.line.fill.background()
add_text_box(slide, 1.7, 1.95, 4, 0.4, '网盘 VIP / SVIP', 14, TEXT, True)
add_text_box(slide, 1.5, 2.5, 5.0, 1.5, '• 存储与传输\n• 影音播放\n• 文件处理\n• AI 智能特权（视频总结/录音纪要/智能整理）', 12, TEXT2)

add_card(slide, 1.2, 4.2, 5.5, 2.0)
shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), Inches(4.4), Inches(0.12), Inches(0.12))
shape.fill.solid(); shape.fill.fore_color.rgb = INDIGO; shape.line.fill.background()
add_text_box(slide, 1.7, 4.35, 4, 0.4, 'AI 会员', 14, TEXT, True)
add_text_box(slide, 1.5, 4.85, 5.0, 1.3, '• 夸克 PPT、文档下载\n• 学习（视频讲题/资源下载）\n• 图像处理、音视频转写\n• 文件格式转换', 12, TEXT2)

# Right column
add_card(slide, 7.0, 1.8, 5.5, 1.3)
shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.3), Inches(2.0), Inches(0.12), Inches(0.12))
shape.fill.solid(); shape.fill.fore_color.rgb = WARM; shape.line.fill.background()
add_text_box(slide, 7.5, 1.95, 4, 0.4, '扫描会员', 14, TEXT, True)
add_text_box(slide, 7.3, 2.45, 5.0, 0.5, '主要用于独立 APP 夸克扫描王', 12, TEXT2)

add_card(slide, 7.0, 3.3, 5.5, 1.3)
shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.3), Inches(3.5), Inches(0.12), Inches(0.12))
shape.fill.solid(); shape.fill.fore_color.rgb = GREEN; shape.line.fill.background()
add_text_box(slide, 7.5, 3.45, 4, 0.4, '小说会员', 14, TEXT, True)
add_text_box(slide, 7.3, 3.95, 5.0, 0.5, '独立内容品类会员', 12, TEXT2)

# Highlight
add_card(slide, 7.0, 4.9, 5.5, 1.3, CYAN)
add_text_box(slide, 7.3, 5.1, 4.9, 0.4, '矛盾', 13, TEXT, True)
add_text_box(slide, 7.3, 5.5, 4.9, 0.6, 'AI 应用集中到夸克了，还是需要开各种会员。一种会员对应一种场景预设，复杂体系削弱了一站式优势。', 12, TEXT2)
add_page_num(slide, 5)

# ═══════════════════════════════════════
# 06 付费对比
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_label(slide, '核心洞察')
add_line(slide, 1.2, 1.1, 0.08)
add_text_box(slide, 1.5, 0.9, 8, 0.7, '付费逻辑的本质差异', 26, TEXT, True)

# vivo card
add_card(slide, 1.2, 1.8, 5.0, 2.5, RGBColor(0xE8,0xC8,0xC8))
add_text_box(slide, 1.5, 2.0, 4.5, 0.5, 'vivo：买空间', 16, RED_S, True)
add_text_box(slide, 1.5, 2.6, 4.5, 1.5, 'VIP 四项权益（超大空间/三重备份/家人共享/高频备份）全部围绕存储可靠性。\n\n用户付费后在功能层面没有任何获得感提升，续费动力弱。', 13, TEXT2)

# VS
add_text_box(slide, 6.25, 2.7, 0.8, 0.5, 'VS', 14, TEXT3, True, PP_ALIGN.CENTER)

# Quark card
add_card(slide, 7.1, 1.8, 5.0, 2.5, RGBColor(0xC8,0xE0,0xE0))
add_text_box(slide, 7.4, 2.0, 4.5, 0.5, '夸克：买能力', 16, RGBColor(0x4A,0x9A,0xAA), True)
add_text_box(slide, 7.4, 2.6, 4.5, 1.5, 'AI 能力、文件处理、影音特权均为会员专属。\n\n用户感知到的是"我能做更多事"，而非"我有更多空间"，留存与续费逻辑更强。', 13, TEXT2)

# Problem box
add_line(slide, 1.2, 4.8, 11, RED_S, 0.04)
add_text_box(slide, 1.5, 5.0, 5, 0.5, 'vivo 的资源闲置问题', 14, TEXT, True)
add_text_box(slide, 1.5, 5.5, 10.5, 1.5,
    '蓝心大模型已具备端云灵活部署能力，意图理解提升约 30%，蓝心小V 已实现通话摘要、图片语义搜索、笔记 AI 写作——但完全未接入云服务增值体系。',
    13, TEXT2)
add_page_num(slide, 6)

# ═══════════════════════════════════════
# 07 核心判断
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_label(slide, '战略判断')
add_text_box(slide, 3, 1.5, 7.3, 0.5, 'CORE INSIGHT', 14, TEXT3, False, PP_ALIGN.CENTER)
add_text_box(slide, 2, 2.3, 9.3, 1.8, 'vivo 云端沉积着\n夸克永远触达不到的数据', 32, TEXT, True, PP_ALIGN.CENTER, '微软雅黑')
add_text_box(slide, 2.5, 4.2, 8.3, 1.0, '通话录音、联系人关系、系统笔记、相册元数据——\n这是 vivo 云服务天然的护城河。',
             16, TEXT2, False, PP_ALIGN.CENTER)
add_text_box(slide, 2.5, 5.3, 8.3, 0.8, '当前这些数据孤立存储，缺乏横向连接与 AI 理解，\n是最值得突破的方向。',
             14, TEXT2, False, PP_ALIGN.CENTER)

tags = ['通话录音', '联系人关系', '系统笔记', '相册元数据']
tag_w = 1.5
start_x = (13.333 - len(tags) * (tag_w + 0.2)) / 2
for i, t in enumerate(tags):
    x = start_x + i * (tag_w + 0.2)
    add_card(slide, x, 6.3, tag_w, 0.45, RGBColor(0xD0,0xDA,0xC8))
    add_text_box(slide, x, 6.35, tag_w, 0.35, t, 11, TEXT2, False, PP_ALIGN.CENTER)
add_page_num(slide, 7)

# ═══════════════════════════════════════
# 08 场景一
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_label(slide, '优化方向')
add_line(slide, 1.2, 1.1, 0.08)
add_text_box(slide, 1.5, 0.9, 8, 0.7, '场景一：个人记忆引擎', 26, TEXT, True)

# Problem
add_card(slide, 1.2, 1.8, 10.9, 1.0)
add_text_box(slide, 1.5, 1.85, 3, 0.4, '问题', 14, TEXT, True)
add_text_box(slide, 1.5, 2.25, 10.3, 0.5, '照片、录音、笔记、联系人分散在不同模块，无法被统一检索和理解。', 13, TEXT2)

# Solution
add_card(slide, 1.2, 3.1, 10.9, 2.8)
add_text_box(slide, 1.5, 3.2, 5, 0.4, '优化方向', 14, CYAN, True)
add_text_box(slide, 1.5, 3.65, 10.3, 0.6, '以蓝心大模型为底座，打通云端私人数据，提供跨模块语义检索能力。', 14, TEXT)
add_multi_text(slide, 1.5, 4.3, 10.3, 1.4, [
    ('→  "找我去年在南京拍的照片" → 调相册元数据 + 地理位置', 13, TEXT, False, PP_ALIGN.LEFT, '微软雅黑'),
    ('→  "整理我和张三的沟通记录" → 关联通话录音摘要 + 聊天笔记', 13, TEXT, False, PP_ALIGN.LEFT, '微软雅黑'),
    ('→  "上周开会说了什么" → 调录音机 + 原子笔记', 13, TEXT, False, PP_ALIGN.LEFT, '微软雅黑'),
])

# Highlight
add_card(slide, 1.2, 6.2, 10.9, 0.9, CYAN)
add_text_box(slide, 1.5, 6.3, 10.3, 0.7, '差异化：基于系统级私人数据，任何第三方云盘（包括夸克）均无法复制。', 14, TEXT, True)
add_page_num(slide, 8)

# ═══════════════════════════════════════
# 09 场景二
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_label(slide, '优化方向')
add_line(slide, 1.2, 1.1, 0.08)
add_text_box(slide, 1.5, 0.9, 8, 0.7, '场景二：系统浏览器 × 云服务', 26, TEXT, True)

# Problem
add_card(slide, 1.2, 1.8, 10.9, 1.0)
add_text_box(slide, 1.5, 1.85, 3, 0.4, '问题', 14, TEXT, True)
add_text_box(slide, 1.5, 2.25, 10.3, 0.5, 'vivo 浏览器作为系统预装应用，用户粘性弱，缺乏主动打开的动力。', 13, TEXT2)

# Solution
add_card(slide, 1.2, 3.1, 10.9, 1.0)
add_text_box(slide, 1.5, 3.2, 5, 0.4, '优化方向', 14, CYAN, True)
add_text_box(slide, 1.5, 3.6, 10.3, 0.5, '在浏览器中植入"搜索我的数据"入口，调用蓝心大模型理解意图，直连云端相册、文件、笔记。', 13, TEXT)

# VS comparison
add_card(slide, 1.2, 4.4, 5.0, 1.5, RGBColor(0xC8,0xE0,0xE0))
add_text_box(slide, 1.5, 4.5, 4.5, 0.4, 'vivo 浏览器搜索', 13, RGBColor(0x4A,0x9A,0xAA), True)
add_text_box(slide, 1.5, 4.95, 4.5, 0.8, '搜索的是用户自己的私人数据——差异化更强，更具隐私价值感', 12, TEXT2)

add_text_box(slide, 6.25, 4.8, 0.8, 0.5, 'VS', 14, TEXT3, True, PP_ALIGN.CENTER)

add_card(slide, 7.1, 4.4, 5.0, 1.5, RGBColor(0xE0,0xD8,0xC8))
add_text_box(slide, 7.4, 4.5, 4.5, 0.4, '夸克 AI 超级框', 13, WARM, True)
add_text_box(slide, 7.4, 4.95, 4.5, 0.8, '搜索的是公网内容', 12, TEXT2)

# Examples
add_card(slide, 1.2, 6.2, 10.9, 0.9)
add_text_box(slide, 1.5, 6.3, 10.3, 0.7,
    '示例交互：  "存的合同" → 云盘文件  |  "老家朋友" → 联系人 + 通话摘要  |  "上次记的菜谱" → 原子笔记',
    13, TEXT2)
add_page_num(slide, 9)

# ═══════════════════════════════════════
# 10 总结
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_label(slide, '总结')
add_text_box(slide, 3, 1.3, 7.3, 0.5, 'NARRATIVE', 14, TEXT3, False, PP_ALIGN.CENTER)
add_text_box(slide, 2, 2.0, 9.3, 1.5, '把分散的私人数据\n变成一个可对话的个人知识库', 28, TEXT, True, PP_ALIGN.CENTER, '微软雅黑')

# Diagram: 3 nodes
nodes = [
    ('蓝心小V', CYAN, 2.5),
    ('vivo 浏览器', INDIGO, 5.2),
    ('个人知识库', GREEN, 7.9),
]
for name, color, x in nodes:
    add_card(slide, x, 3.8, 2.2, 0.7, color)
    add_text_box(slide, x, 3.9, 2.2, 0.5, name, 13, TEXT, True, PP_ALIGN.CENTER)

add_text_box(slide, 4.7, 3.9, 0.5, 0.5, '+', 20, TEXT3, False, PP_ALIGN.CENTER)
add_text_box(slide, 7.4, 3.9, 0.5, 0.5, '→', 20, TEXT3, False, PP_ALIGN.CENTER)

add_multi_text(slide, 2.5, 4.8, 8.3, 1.8, [
    ('两个入口：蓝心小V 和 vivo 浏览器', 14, TEXT, False, PP_ALIGN.CENTER, '微软雅黑'),
    ('内容层：云端私人数据', 14, TEXT, False, PP_ALIGN.CENTER, '微软雅黑'),
    ('理解层：蓝心大模型', 14, TEXT, False, PP_ALIGN.CENTER, '微软雅黑'),
    ('三者打通，形成夸克无法复制的差异化能力。', 14, TEXT2, False, PP_ALIGN.CENTER, '微软雅黑'),
])

add_line(slide, 6.1, 6.5, 1)
add_text_box(slide, 3, 6.7, 7.3, 0.5, '精简功能设计 · 一站式一体化 · 系统级优势',
             13, TEXT3, False, PP_ALIGN.CENTER)
add_page_num(slide, 10)

# Save
out = os.path.join(os.path.dirname(__file__), 'output', 'vivo-cloud-analysis.pptx')
prs.save(out)
print(f'Saved: {out}')
